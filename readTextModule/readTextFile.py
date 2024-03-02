# function for read pdf and return a text
from PyPDF2 import PdfReader
# function for read docx
from docx import Document
# function for read pptx
from pptx import Presentation


def read_pdf(file_path):
    with open(file_path, 'rb') as file:
        reader = PdfReader(file)
        text = ""
        for page_num in range(reader.pages.__len__()):
            text += reader.pages[page_num].extract_text().replace("\n", " ")
    return text


def read_docx(file_path):
    doc = Document(file_path)
    text = ""
    for paragraph in doc.paragraphs:
        text += paragraph.text + "\n"
    return text


def read_ppt(file_path):
    presentation = Presentation(file_path)
    text = ""
    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text


